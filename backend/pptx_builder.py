import json
import re
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def parse_slide_content(raw: str) -> tuple[str, list[str]]:
    """Extract title and bullets from the JSON produced by the PPTX aggregate step."""
    stripped = raw.strip()
    # Strip markdown code fences if present
    stripped = re.sub(r"^```json\s*", "", stripped)
    stripped = re.sub(r"```\s*$", "", stripped)
    stripped = stripped.strip()
    try:
        data = json.loads(stripped)
        title = data.get("title", "")
        bullets = data.get("bullets", [])
        if isinstance(bullets, list):
            return title, [str(b) for b in bullets]
    except (json.JSONDecodeError, TypeError):
        pass
    # Fallback: treat first line as title, remaining non-empty lines as bullets
    lines = [l.strip() for l in stripped.splitlines() if l.strip()]
    title = lines[0].lstrip("#").strip() if lines else ""
    bullets = [l.lstrip("-•* ").strip() for l in lines[1:] if l.strip()]
    return title, bullets


def build_slide(raw_content: str, slide_config, chapter_title: str) -> bytes:
    """Build a single-slide PPTX from the raw agent output and SlideConfig settings."""
    title_text, bullets = parse_slide_content(raw_content)
    if not title_text:
        title_text = chapter_title

    prs = Presentation()

    # Set slide dimensions
    if slide_config.slide_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:  # 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(slide_config.bg_color)

    w = prs.slide_width
    h = prs.slide_height
    ml = Inches(slide_config.margin_left)
    mt = Inches(slide_config.margin_top)
    mr = Inches(slide_config.margin_right)
    mb = Inches(slide_config.margin_bottom)
    content_w = w - ml - mr

    # Title text box — top ~25% of content area
    title_h = Inches(1.4)
    title_box = slide.shapes.add_textbox(ml, mt, content_w, title_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.name = slide_config.title_font
    run.font.size = Pt(slide_config.title_font_size)
    run.font.bold = slide_config.title_bold
    run.font.color.rgb = _hex_to_rgb(slide_config.title_color)

    # Body text box — remaining space
    body_top = mt + title_h + Inches(0.15)
    body_h = h - body_top - mb
    body_box = slide.shapes.add_textbox(ml, body_top, content_w, body_h)
    btf = body_box.text_frame
    btf.word_wrap = True

    for i, bullet in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = f"• {bullet}"
        run = para.runs[0]
        run.font.name = slide_config.body_font
        run.font.size = Pt(slide_config.body_font_size)
        run.font.bold = slide_config.body_bold
        run.font.color.rgb = _hex_to_rgb(slide_config.body_color)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
